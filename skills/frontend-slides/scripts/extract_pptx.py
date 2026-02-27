#!/usr/bin/env python3
"""
Extract content from a PowerPoint file into a JSON structure.

Usage:
    python extract_pptx.py <input.pptx> <output_dir>

Output:
    <output_dir>/slides.json   — structured slide content
    <output_dir>/assets/       — extracted images

Exits with code 1 on error and prints a human-readable message to stderr.
"""

import sys
import os
import json


def check_dependency():
    try:
        import pptx  # noqa: F401
    except ImportError:
        print(
            "Error: python-pptx is not installed.\n"
            "Install it with:  pip install python-pptx",
            file=sys.stderr,
        )
        sys.exit(1)


def extract_pptx(file_path, output_dir):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    if not os.path.isfile(file_path):
        print(f"Error: file not found: {file_path}", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(file_path)
    slides_data = []

    assets_dir = os.path.join(output_dir, "assets")
    os.makedirs(assets_dir, exist_ok=True)

    for slide_num, slide in enumerate(prs.slides):
        slide_data = {
            "number": slide_num + 1,
            "title": "",
            "content": [],
            "images": [],
            "notes": "",
            "unsupported": [],
        }

        for shape in slide.shapes:
            # Title
            if shape.has_text_frame and shape == slide.shapes.title:
                slide_data["title"] = shape.text_frame.text.strip()
                continue

            # Text blocks
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    slide_data["content"].append({"type": "text", "content": text})
                continue

            # Pictures
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                ext = image.ext
                name = f"slide{slide_num + 1}_img{len(slide_data['images']) + 1}.{ext}"
                path = os.path.join(assets_dir, name)
                with open(path, "wb") as f:
                    f.write(image.blob)
                slide_data["images"].append(
                    {
                        "path": f"assets/{name}",
                        "width": shape.width,
                        "height": shape.height,
                    }
                )
                continue

            # Tables
            if shape.has_table:
                rows = []
                for row in shape.table.rows:
                    rows.append([cell.text.strip() for cell in row.cells])
                slide_data["content"].append({"type": "table", "rows": rows})
                continue

            # Charts, grouped shapes, media, etc. — log but don't fail
            shape_type_name = str(shape.shape_type)
            slide_data["unsupported"].append(
                {"shape_type": shape_type_name, "name": shape.name}
            )

        # Speaker notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_data["notes"] = notes_text

        slides_data.append(slide_data)

    output_path = os.path.join(output_dir, "slides.json")
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(slides_data, f, indent=2, ensure_ascii=False)

    # Summary to stdout for the agent to read
    total_images = sum(len(s["images"]) for s in slides_data)
    unsupported_count = sum(len(s["unsupported"]) for s in slides_data)
    print(json.dumps({
        "slide_count": len(slides_data),
        "total_images": total_images,
        "unsupported_shapes": unsupported_count,
        "output_json": output_path,
        "assets_dir": assets_dir,
    }))


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print(f"Usage: {sys.argv[0]} <input.pptx> <output_dir>", file=sys.stderr)
        sys.exit(1)

    check_dependency()
    extract_pptx(sys.argv[1], sys.argv[2])
